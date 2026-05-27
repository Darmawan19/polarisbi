"""
deck_design.py — PolarisBI PPTX Design System

Design tokens, helpers, and base components adapted from CX Intelligence.
Slide builders will be added in Day 3 (see TODO block below).
"""

import io

import matplotlib
import matplotlib.pyplot as plt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── DESIGN TOKENS ────────────────────────────────────────────────────────────
# Per research: white bg, gray base, single accent, semantic colors used sparingly

# Colors
BG_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY    = RGBColor(0x1F, 0x29, 0x37)   # near-black
TEXT_SECONDARY  = RGBColor(0x6B, 0x72, 0x80)   # muted gray
GRID_LIGHT      = RGBColor(0xE5, 0xE7, 0xEB)
BASE_GRAY       = RGBColor(0xBF, 0xBF, 0xBF)   # Cole's base
ACCENT          = RGBColor(0x1F, 0x4E, 0x79)   # consulting blue
ACCENT_LIGHT    = RGBColor(0xDB, 0xE4, 0xF0)   # lightened accent for fills
SEMANTIC_RED    = RGBColor(0xC0, 0x39, 0x2B)
SEMANTIC_GREEN  = RGBColor(0x27, 0xAE, 0x60)
SEMANTIC_AMBER  = RGBColor(0xD6, 0x89, 0x10)

# Matplotlib hex equivalents
MPL_ACCENT      = "#1F4E79"
MPL_BASE_GRAY   = "#BFBFBF"
MPL_TEXT        = "#1F2937"
MPL_GRID        = "#E5E7EB"
MPL_AMBER       = "#D68910"
MPL_RED         = "#C0392B"
MPL_GREEN       = "#27AE60"

# Typography (system-safe substitutes for proprietary consulting fonts)
FONT_HEADLINE   = "Georgia"   # McKinsey-style serif headlines
FONT_BODY       = "Arial"     # Universal body

# Sizes
SIZE_TITLE      = Pt(34)      # cover slide
SIZE_SUBTITLE   = Pt(18)
SIZE_HEADLINE   = Pt(22)      # slide action titles
SIZE_KPI_VALUE  = Pt(40)
SIZE_BODY       = Pt(14)
SIZE_LABEL      = Pt(11)
SIZE_FOOTNOTE   = Pt(9)


# ── HELPERS ──────────────────────────────────────────────────────────────────

def set_run(run, text, *, font=FONT_BODY, size=SIZE_BODY, bold=False,
            italic=False, color=TEXT_PRIMARY):
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, *, font=FONT_BODY, size=SIZE_BODY,
             bold=False, italic=False, color=TEXT_PRIMARY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, font=font, size=size, bold=bold, italic=italic, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x, y, w, h, color, weight=1.0):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y),
                                      Inches(x + w), Inches(y + h))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def slide_chrome(slide, headline, subhead=None, page_num=None, total_pages=None):
    """Apply consistent header to every content slide."""
    bg = add_rect(slide, 0, 0, 13.333, 7.5, BG_WHITE)
    bg.shadow.inherit = False

    add_text(slide, headline, 0.6, 0.45, 12.1, 0.7,
             font=FONT_HEADLINE, size=SIZE_HEADLINE, bold=True,
             color=TEXT_PRIMARY)

    if subhead:
        add_text(slide, subhead, 0.6, 1.15, 12.1, 0.35,
                 size=Pt(12), color=TEXT_SECONDARY, italic=True)

    add_line(slide, 0.6, 1.55, 12.1, 0, GRID_LIGHT, 0.75)

    if page_num and total_pages:
        add_text(slide, f"{page_num} / {total_pages}",
                 11.5, 7.05, 1.2, 0.3,
                 size=SIZE_FOOTNOTE, color=TEXT_SECONDARY,
                 align=PP_ALIGN.RIGHT)

    add_text(slide, "PolarisBI  ·  AI Cockpit untuk IT BA Asuransi",
             0.6, 7.05, 8, 0.3, size=SIZE_FOOTNOTE, color=TEXT_SECONDARY)


def configure_matplotlib():
    """Apply consulting-style defaults to matplotlib."""
    matplotlib.rcParams.update({
        'figure.dpi': 200,
        'savefig.dpi': 200,
        'savefig.bbox': 'tight',
        'savefig.facecolor': 'white',
        'font.family': 'DejaVu Sans',
        'font.size': 11,
        'axes.titlesize': 13,
        'axes.titleweight': 'bold',
        'axes.titlecolor': MPL_TEXT,
        'axes.labelsize': 11,
        'axes.labelcolor': MPL_TEXT,
        'axes.edgecolor': '#CBD5E1',
        'axes.linewidth': 0.75,
        'axes.spines.top': False,
        'axes.spines.right': False,
        'axes.grid': True,
        'axes.axisbelow': True,
        'grid.color': MPL_GRID,
        'grid.linewidth': 0.5,
        'grid.alpha': 0.8,
        'xtick.color': MPL_TEXT,
        'ytick.color': MPL_TEXT,
        'xtick.labelsize': 10,
        'ytick.labelsize': 10,
        'legend.frameon': False,
        'legend.fontsize': 10,
    })


def chart_to_image(fig):
    """Convert matplotlib fig to image bytes for python-pptx insertion."""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


# ── COMPONENTS ───────────────────────────────────────────────────────────────

def kpi_tile(slide, x, y, w, h, value, label, delta=None, delta_color=None):
    """A single KPI tile: big number, label, optional delta arrow."""
    add_rect(slide, x, y, w, h, BG_WHITE, GRID_LIGHT, 0.75)

    add_text(slide, str(value), x + 0.15, y + 0.18, w - 0.3, h * 0.45,
             font=FONT_HEADLINE, size=SIZE_KPI_VALUE, bold=True,
             color=TEXT_PRIMARY, align=PP_ALIGN.LEFT)

    add_text(slide, label, x + 0.15, y + h * 0.55, w - 0.3, h * 0.22,
             size=Pt(11), color=TEXT_SECONDARY, align=PP_ALIGN.LEFT)

    if delta:
        color = delta_color if delta_color else TEXT_SECONDARY
        add_text(slide, delta, x + 0.15, y + h * 0.78, w - 0.3, h * 0.2,
                 size=Pt(11), bold=True, color=color, align=PP_ALIGN.LEFT)


def quote_callout(slide, x, y, w, h, quote, attribution):
    """Verbatim quote callout — italic body, attribution underneath."""
    add_rect(slide, x, y, w, h, RGBColor(0xF9, 0xFA, 0xFB), GRID_LIGHT, 0.5)

    add_text(slide, "“", x + 0.1, y + 0.05, 0.5, 0.5,
             font=FONT_HEADLINE, size=Pt(36), bold=True, color=ACCENT)

    add_text(slide, quote, x + 0.55, y + 0.15, w - 0.7, h - 0.6,
             size=Pt(11), italic=True, color=TEXT_PRIMARY)

    add_text(slide, f"— {attribution}", x + 0.55, y + h - 0.35, w - 0.7, 0.25,
             size=Pt(9), color=TEXT_SECONDARY)


# ─── TODO: Slide Builders (Rabu pagi) ─────────────────────────────────────
# build_query_result_slide(question, sql, table, chart) — slide hasil query
# build_kpi_dashboard_slide(kpis) — snapshot KPI insurance
# build_executive_brief(findings, charts) — morning brief format
