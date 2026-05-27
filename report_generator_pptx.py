"""
report_generator_pptx.py — Step 8 of CX Intelligence Pipeline (v2 — consulting-grade)

Generates a 7-slide monthly CX research deck following consulting/PM best practices:
- White background, gray-base + single-accent palette (McKinsey-inspired)
- Action titles (full sentences with insight + number)
- Pyramid Principle / SCR executive summary
- Verbatim user quote callouts
- KPI tiles with delta arrows
- Watch Area slide (signal -> impact -> trigger)
- Recommendations slide (action + owner + timeline + metric)
- Charts via matplotlib (PNG embed), no pies, no 3D, no defaults
"""

import io
import os
from datetime import date

import matplotlib.pyplot as plt
import matplotlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── DESIGN TOKENS ────────────────────────────────────────────────────────────
# Per research: white bg, gray base, single accent, semantic colors used sparingly

# Colors
BG_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_PRIMARY    = RGBColor(0x1F, 0x29, 0x37)   # near-black
TEXT_SECONDARY  = RGBColor(0x6B, 0x72, 0x80)   # muted gray
GRID_LIGHT      = RGBColor(0xE5, 0xE7, 0xEB)
BASE_GRAY       = RGBColor(0xBF, 0xBF, 0xBF)   # Cole's base
ACCENT          = RGBColor(0x1F, 0x4E, 0x79)   # consulting blue
ACCENT_LIGHT    = RGBColor(0xDB, 0xE4, 0xF0)   # lightened accent for fills
SEMANTIC_RED    = RGBColor(0xC0, 0x39, 0x2B)
SEMANTIC_GREEN  = RGBColor(0x27, 0xAE, 0x60)
SEMANTIC_AMBER  = RGBColor(0xD6, 0x89, 0x10)

# Matplotlib hex equivalents
MPL_ACCENT      = "#1F4E79"
MPL_BASE_GRAY   = "#BFBFBF"
MPL_TEXT        = "#1F2937"
MPL_GRID        = "#E5E7EB"
MPL_AMBER       = "#D68910"
MPL_RED         = "#C0392B"
MPL_GREEN       = "#27AE60"

# Typography (system-safe substitutes for proprietary consulting fonts)
FONT_HEADLINE   = "Georgia"   # McKinsey-style serif headlines
FONT_BODY       = "Arial"     # Universal body

# Sizes
SIZE_TITLE      = Pt(34)      # cover slide
SIZE_SUBTITLE   = Pt(18)
SIZE_HEADLINE   = Pt(22)      # slide action titles
SIZE_KPI_VALUE  = Pt(40)
SIZE_BODY       = Pt(14)
SIZE_LABEL      = Pt(11)
SIZE_FOOTNOTE   = Pt(9)


# ── HELPERS ──────────────────────────────────────────────────────────────────

def set_run(run, text, *, font=FONT_BODY, size=SIZE_BODY, bold=False,
            italic=False, color=TEXT_PRIMARY):
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, *, font=FONT_BODY, size=SIZE_BODY,
             bold=False, italic=False, color=TEXT_PRIMARY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, font=font, size=size, bold=bold, italic=italic, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x, y, w, h, color, weight=1.0):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y),
                                      Inches(x + w), Inches(y + h))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def slide_chrome(slide, headline, subhead=None, page_num=None, total_pages=None):
    """Apply consistent header to every content slide."""
    # White background
    bg = add_rect(slide, 0, 0, 13.333, 7.5, BG_WHITE)
    bg.shadow.inherit = False

    # Action title (headline) - the most important element
    add_text(slide, headline, 0.6, 0.45, 12.1, 0.7,
             font=FONT_HEADLINE, size=SIZE_HEADLINE, bold=True,
             color=TEXT_PRIMARY)

    if subhead:
        add_text(slide, subhead, 0.6, 1.15, 12.1, 0.35,
                 size=Pt(12), color=TEXT_SECONDARY, italic=True)

    # Thin divider under header
    add_line(slide, 0.6, 1.55, 12.1, 0, GRID_LIGHT, 0.75)

    # Page number bottom-right
    if page_num and total_pages:
        add_text(slide, f"{page_num} / {total_pages}",
                 11.5, 7.05, 1.2, 0.3,
                 size=SIZE_FOOTNOTE, color=TEXT_SECONDARY,
                 align=PP_ALIGN.RIGHT)

    # Brand footer bottom-left
    add_text(slide, "CX Intelligence  ·  github.com/Darmawan19/cx-intelligence",
             0.6, 7.05, 8, 0.3, size=SIZE_FOOTNOTE, color=TEXT_SECONDARY)


def configure_matplotlib():
    """Apply consulting-style defaults to matplotlib."""
    matplotlib.rcParams.update({
        'figure.dpi': 200,
        'savefig.dpi': 200,
        'savefig.bbox': 'tight',
        'savefig.facecolor': 'white',
        'font.family': 'DejaVu Sans',  # Arial-equivalent that's reliably available
        'font.size': 11,
        'axes.titlesize': 13,
        'axes.titleweight': 'bold',
        'axes.titlecolor': MPL_TEXT,
        'axes.labelsize': 11,
        'axes.labelcolor': MPL_TEXT,
        'axes.edgecolor': '#CBD5E1',
        'axes.linewidth': 0.75,
        'axes.spines.top': False,
        'axes.spines.right': False,
        'axes.grid': True,
        'axes.axisbelow': True,
        'grid.color': MPL_GRID,
        'grid.linewidth': 0.5,
        'grid.alpha': 0.8,
        'xtick.color': MPL_TEXT,
        'ytick.color': MPL_TEXT,
        'xtick.labelsize': 10,
        'ytick.labelsize': 10,
        'legend.frameon': False,
        'legend.fontsize': 10,
    })


def chart_to_image(fig):
    """Convert matplotlib fig to image bytes for python-pptx insertion."""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


# ── COMPONENTS ───────────────────────────────────────────────────────────────

def kpi_tile(slide, x, y, w, h, value, label, delta=None, delta_color=None):
    """A single KPI tile: big number, label, optional delta arrow."""
    # Thin border, white fill
    border = add_rect(slide, x, y, w, h, BG_WHITE, GRID_LIGHT, 0.75)

    # Large value
    add_text(slide, str(value), x + 0.15, y + 0.18, w - 0.3, h * 0.45,
             font=FONT_HEADLINE, size=SIZE_KPI_VALUE, bold=True,
             color=TEXT_PRIMARY, align=PP_ALIGN.LEFT)

    # Label below
    add_text(slide, label, x + 0.15, y + h * 0.55, w - 0.3, h * 0.22,
             size=Pt(11), color=TEXT_SECONDARY, align=PP_ALIGN.LEFT)

    # Delta indicator
    if delta:
        color = delta_color if delta_color else TEXT_SECONDARY
        add_text(slide, delta, x + 0.15, y + h * 0.78, w - 0.3, h * 0.2,
                 size=Pt(11), bold=True, color=color, align=PP_ALIGN.LEFT)


def quote_callout(slide, x, y, w, h, quote, attribution):
    """Verbatim user quote — italic body, attribution underneath."""
    add_rect(slide, x, y, w, h, RGBColor(0xF9, 0xFA, 0xFB), GRID_LIGHT, 0.5)

    # Opening quote mark in accent
    add_text(slide, "\u201C", x + 0.1, y + 0.05, 0.5, 0.5,
             font=FONT_HEADLINE, size=Pt(36), bold=True, color=ACCENT)

    # Quote body
    add_text(slide, quote, x + 0.55, y + 0.15, w - 0.7, h - 0.6,
             size=Pt(11), italic=True, color=TEXT_PRIMARY)

    # Attribution at bottom
    add_text(slide, f"— {attribution}", x + 0.55, y + h - 0.35, w - 0.7, 0.25,
             size=Pt(9), color=TEXT_SECONDARY)


# ── CHART BUILDERS ───────────────────────────────────────────────────────────

def priority_bar_chart(items, emphasize_top_n=3, watch_indices=None,
                       figsize=(6.5, 4.0)):
    """Horizontal bar chart: feature area priority scores.

    items: list of (label, score) tuples, already sorted desc
    emphasize_top_n: top N get accent color
    watch_indices: list of indices to color amber
    """
    watch_indices = watch_indices or []
    labels = [it[0] for it in items]
    scores = [it[1] for it in items]
    n = len(items)

    colors = []
    for i, _ in enumerate(items):
        if i < emphasize_top_n:
            colors.append(MPL_ACCENT)
        elif i in watch_indices:
            colors.append(MPL_AMBER)
        else:
            colors.append(MPL_BASE_GRAY)

    fig, ax = plt.subplots(figsize=figsize)
    y_pos = range(n)
    bars = ax.barh(y_pos, scores, color=colors, height=0.65)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=11)
    ax.invert_yaxis()
    ax.set_xlabel("Priority Score  (Volume × Avg Severity²)",
                  fontsize=10, color=MPL_TEXT)

    # Direct value labels at end of each bar
    max_score = max(scores)
    for i, (bar, score) in enumerate(zip(bars, scores)):
        label_text = f"{score:.0f}"
        if i < emphasize_top_n or i in watch_indices:
            weight = 'bold'
            color = colors[i]
        else:
            weight = 'normal'
            color = MPL_TEXT
        ax.text(bar.get_width() + max_score * 0.015,
                bar.get_y() + bar.get_height() / 2,
                label_text, va='center', fontsize=10,
                fontweight=weight, color=color)

    ax.set_xlim(0, max_score * 1.15)
    ax.spines['left'].set_visible(False)
    ax.grid(axis='x', alpha=0.4)
    ax.tick_params(axis='y', length=0)
    fig.tight_layout()
    return fig


def competitive_grouped_bar(categories, shopee_vals, tokopedia_vals,
                             figsize=(6.5, 3.8)):
    """Grouped bar: Shopee vs Tokopedia normalized scores per category."""
    fig, ax = plt.subplots(figsize=figsize)
    x = range(len(categories))
    width = 0.36

    bars_s = ax.bar([i - width/2 for i in x], shopee_vals, width,
                    label="Shopee", color=MPL_ACCENT, edgecolor='none')
    bars_t = ax.bar([i + width/2 for i in x], tokopedia_vals, width,
                    label="Tokopedia", color=MPL_BASE_GRAY, edgecolor='none')

    # Value labels on top
    for bars in (bars_s, bars_t):
        for bar in bars:
            h = bar.get_height()
            ax.text(bar.get_x() + bar.get_width() / 2, h + 3,
                    f"{h:.0f}", ha='center', va='bottom',
                    fontsize=9, color=MPL_TEXT)

    ax.set_xticks(list(x))
    ax.set_xticklabels(categories, fontsize=10)
    ax.set_ylabel("Score per 100 reviews", fontsize=10, color=MPL_TEXT)
    ax.legend(loc='upper right', fontsize=10)
    ax.spines['left'].set_visible(False)
    ax.grid(axis='y', alpha=0.4)
    ax.tick_params(axis='x', length=0)
    fig.tight_layout()
    return fig


def mom_change_chart(items, figsize=(5.5, 3.5)):
    """Horizontal lollipop chart: % MoM change per area."""
    labels = [it[0] for it in items]
    changes = [it[1] for it in items]
    n = len(items)

    colors = [MPL_AMBER if c > 50 else MPL_GREEN if c < 0 else MPL_ACCENT
              for c in changes]

    fig, ax = plt.subplots(figsize=figsize)
    y_pos = list(range(n))

    for i, (label, change, color) in enumerate(zip(labels, changes, colors)):
        ax.hlines(y=i, xmin=0, xmax=change, color=color, linewidth=2)
        ax.scatter(change, i, s=120, color=color, zorder=3, edgecolor='white',
                   linewidth=2)
        # Value at end
        offset = 3 if change >= 0 else -3
        ha = 'left' if change >= 0 else 'right'
        ax.text(change + offset, i,
                f"{'+' if change > 0 else ''}{change:.0f}%",
                va='center', ha=ha, fontsize=10, fontweight='bold',
                color=color)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=10)
    ax.invert_yaxis()
    ax.axvline(0, color=MPL_BASE_GRAY, linewidth=0.75)
    ax.set_xlabel("Month-over-Month Change in Priority Score",
                  fontsize=10, color=MPL_TEXT)
    ax.set_xlim(min(changes) * 1.25, max(changes) * 1.25)
    ax.spines['left'].set_visible(False)
    ax.grid(axis='x', alpha=0.3)
    ax.tick_params(axis='y', length=0)
    fig.tight_layout()
    return fig


# ── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def build_cover(prs, *, app_name, period, run_label, total_reviews,
                pct_negative, avg_severity, top_issue, author):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_rect(s, 0, 0, 13.333, 7.5, BG_WHITE)

    # Thin accent rule top-left
    add_rect(s, 0.6, 0.6, 1.5, 0.06, ACCENT)

    # Eyebrow label
    add_text(s, "MONTHLY CONSUMER LISTENING REPORT",
             0.6, 0.85, 12, 0.3, size=Pt(11), bold=True,
             color=ACCENT)

    # Main title — large serif
    add_text(s, "CX Intelligence",
             0.6, 1.5, 12, 1.1,
             font=FONT_HEADLINE, size=Pt(54), bold=True,
             color=TEXT_PRIMARY)

    # Subtitle — descriptive sentence
    add_text(s, f"{app_name} · {period} · {run_label}",
             0.6, 2.85, 12, 0.5,
             size=Pt(18), color=TEXT_SECONDARY)

    # Headline finding
    add_text(s, f"This month's top issue: {top_issue}",
             0.6, 3.8, 12, 0.45,
             font=FONT_HEADLINE, size=Pt(18), italic=True,
             color=TEXT_PRIMARY)

    # Key facts row
    facts_y = 4.7
    add_rect(s, 0.6, facts_y - 0.05, 12.1, 0.9,
             RGBColor(0xF9, 0xFA, 0xFB), GRID_LIGHT, 0.5)
    facts = [
        ("Reviews analyzed", str(total_reviews)),
        ("Negative sentiment", f"{pct_negative:.1f}%"),
        ("Avg severity", f"{avg_severity:.2f} / 5"),
    ]
    fact_w = 12.1 / 3
    for i, (label, value) in enumerate(facts):
        fx = 0.6 + i * fact_w
        add_text(s, label, fx + 0.25, facts_y + 0.05, fact_w - 0.4, 0.25,
                 size=Pt(10), color=TEXT_SECONDARY)
        add_text(s, value, fx + 0.25, facts_y + 0.3, fact_w - 0.4, 0.5,
                 font=FONT_HEADLINE, size=Pt(24), bold=True,
                 color=TEXT_PRIMARY)

    # Author / date footer
    add_text(s, f"{author}  ·  {date.today().strftime('%B %Y')}",
             0.6, 7.1, 12, 0.3, size=Pt(10), color=TEXT_SECONDARY)


def build_exec_summary(prs, *, headline, situation, complication,
                        resolution_items, page=2, total=7):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_chrome(s, headline, page_num=page, total_pages=total)

    # SCR Layout: vertical stack of three sections
    y = 1.85
    section_gap = 0.15

    # Section: SITUATION
    add_text(s, "SITUATION", 0.6, y, 1.6, 0.3,
             size=Pt(10), bold=True, color=ACCENT)
    add_text(s, situation, 2.3, y - 0.02, 10.4, 0.8,
             size=Pt(13), color=TEXT_PRIMARY)

    y += 0.85 + section_gap
    add_line(s, 0.6, y - 0.1, 12.1, 0, GRID_LIGHT, 0.5)

    # Section: COMPLICATION
    add_text(s, "COMPLICATION", 0.6, y, 1.6, 0.3,
             size=Pt(10), bold=True, color=ACCENT)
    add_text(s, complication, 2.3, y - 0.02, 10.4, 0.8,
             size=Pt(13), color=TEXT_PRIMARY)

    y += 0.85 + section_gap
    add_line(s, 0.6, y - 0.1, 12.1, 0, GRID_LIGHT, 0.5)

    # Section: RESOLUTION (numbered bold + supporting)
    add_text(s, "RESOLUTION", 0.6, y, 1.6, 0.3,
             size=Pt(10), bold=True, color=ACCENT)

    yy = y - 0.02
    for i, (bold_claim, supporting) in enumerate(resolution_items, 1):
        # Numbered claim
        num_box = add_rect(s, 2.3, yy + 0.05, 0.32, 0.32, ACCENT)
        add_text(s, str(i), 2.3, yy + 0.05, 0.32, 0.32,
                 font=FONT_BODY, size=Pt(11), bold=True,
                 color=BG_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

        add_text(s, bold_claim, 2.75, yy, 9.9, 0.32,
                 size=Pt(12), bold=True, color=TEXT_PRIMARY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, supporting, 2.75, yy + 0.32, 9.9, 0.4,
                 size=Pt(10), color=TEXT_SECONDARY)
        yy += 0.85


def build_headline_kpis(prs, *, headline, kpis, page=3, total=7):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_chrome(s, headline, page_num=page, total_pages=total)

    # 4 KPI tiles in a row
    n = len(kpis)
    tile_h = 2.4
    gap = 0.2
    total_w = 12.1
    tile_w = (total_w - gap * (n - 1)) / n

    y = 2.0
    for i, kpi in enumerate(kpis):
        x = 0.6 + i * (tile_w + gap)
        kpi_tile(s, x, y, tile_w, tile_h,
                 kpi['value'], kpi['label'],
                 kpi.get('delta'), kpi.get('delta_color'))

    # Insight callout below
    insight_y = y + tile_h + 0.4
    add_rect(s, 0.6, insight_y, 12.1, 1.0,
             RGBColor(0xF9, 0xFA, 0xFB), ACCENT, 0)
    add_rect(s, 0.6, insight_y, 0.08, 1.0, ACCENT)  # left accent bar
    add_text(s, "READ:", 0.85, insight_y + 0.15, 1.0, 0.3,
             size=Pt(10), bold=True, color=ACCENT)
    insight_text = ("Negative sentiment held above 75% for the third consecutive month. "
                    "Avg severity ticked up to 2.67/5 — driven by escalating Return & "
                    "Refund issues. Delivery remains the dominant pain category.")
    add_text(s, insight_text, 1.7, insight_y + 0.12, 10.9, 0.85,
             size=Pt(11), color=TEXT_PRIMARY)


def build_deep_dive(prs, *, headline, subhead, chart_image, quote_text,
                     quote_attribution, root_cause_items, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_chrome(s, headline, subhead=subhead, page_num=page, total_pages=total)

    # Layout: chart left (60%), root cause + quote right (40%)
    content_y = 1.85
    content_h = 4.9

    # Chart (left)
    chart_w = 7.0
    s.shapes.add_picture(chart_image, Inches(0.6), Inches(content_y),
                         width=Inches(chart_w), height=Inches(content_h * 0.85))

    # Right column
    right_x = 0.6 + chart_w + 0.3
    right_w = 12.7 - right_x

    # Root cause section
    add_text(s, "ROOT CAUSE HYPOTHESIS", right_x, content_y, right_w, 0.25,
             size=Pt(10), bold=True, color=ACCENT)
    yy = content_y + 0.32
    for i, (claim, support) in enumerate(root_cause_items, 1):
        # Number + claim
        add_text(s, f"{i}.  {claim}", right_x, yy, right_w, 0.4,
                 size=Pt(11), bold=True, color=TEXT_PRIMARY)
        add_text(s, support, right_x + 0.2, yy + 0.4, right_w - 0.2, 0.55,
                 size=Pt(10), color=TEXT_SECONDARY)
        yy += 1.05

    # Verbatim quote at bottom
    quote_y = content_y + content_h - 1.05
    quote_callout(s, 0.6, quote_y, chart_w, 1.0, quote_text, quote_attribution)


def build_competitive(prs, *, headline, subhead, chart_image,
                       finding_items, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_chrome(s, headline, subhead=subhead, page_num=page, total_pages=total)

    # Chart top
    chart_y = 1.9
    chart_w = 12.1
    s.shapes.add_picture(chart_image, Inches(0.6), Inches(chart_y),
                         width=Inches(chart_w), height=Inches(3.3))

    # 3 finding cards below
    cards_y = chart_y + 3.4
    n = len(finding_items)
    card_w = (12.1 - 0.3 * (n - 1)) / n
    card_h = 1.3

    for i, (title, body, color) in enumerate(finding_items):
        cx = 0.6 + i * (card_w + 0.3)
        add_rect(s, cx, cards_y, card_w, card_h, BG_WHITE, GRID_LIGHT, 0.5)
        add_rect(s, cx, cards_y, card_w, 0.06, color)  # top accent
        add_text(s, title, cx + 0.2, cards_y + 0.18, card_w - 0.4, 0.35,
                 size=Pt(11), bold=True, color=TEXT_PRIMARY)
        add_text(s, body, cx + 0.2, cards_y + 0.55, card_w - 0.4, 0.7,
                 size=Pt(10), color=TEXT_SECONDARY)


def build_watch_area(prs, *, headline, subhead, watch_items, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_chrome(s, headline, subhead=subhead, page_num=page, total_pages=total)

    # 3-column layout: Signal | Hypothesized Impact | Monitoring Trigger
    y = 1.95
    col_headers = ["EARLY SIGNAL", "HYPOTHESIZED IMPACT", "MONITORING TRIGGER"]
    col_w = 12.1 / 3
    for i, header in enumerate(col_headers):
        cx = 0.6 + i * col_w
        add_text(s, header, cx + 0.2, y, col_w - 0.4, 0.3,
                 size=Pt(10), bold=True, color=ACCENT)

    add_line(s, 0.6, y + 0.4, 12.1, 0, GRID_LIGHT, 0.5)

    # Rows
    row_h = 1.3
    yy = y + 0.55
    for item in watch_items:
        # Item name header spanning all 3 columns
        add_rect(s, 0.6, yy, 0.08, row_h - 0.2, SEMANTIC_AMBER)
        add_text(s, item['name'], 0.85, yy, 12, 0.3,
                 size=Pt(12), bold=True, color=SEMANTIC_AMBER)
        # Three columns
        for i, key in enumerate(['signal', 'impact', 'trigger']):
            cx = 0.85 + i * col_w
            add_text(s, item[key], cx, yy + 0.35, col_w - 0.4, row_h - 0.45,
                     size=Pt(10), color=TEXT_PRIMARY)
        yy += row_h
        add_line(s, 0.6, yy - 0.1, 12.1, 0, GRID_LIGHT, 0.3)


def build_recommendations(prs, *, headline, subhead, recommendations,
                            page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_chrome(s, headline, subhead=subhead, page_num=page, total_pages=total)

    # Column headers
    y = 1.9
    headers = [("Recommendation", 5.4),
               ("Owner", 2.2),
               ("Timeline", 1.6),
               ("Success Metric", 2.9)]

    xx = 0.6
    for header, w in headers:
        add_text(s, header.upper(), xx + 0.1, y, w - 0.2, 0.3,
                 size=Pt(10), bold=True, color=ACCENT)
        xx += w

    add_line(s, 0.6, y + 0.35, 12.1, 0, GRID_LIGHT, 0.5)

    # Rows
    yy = y + 0.55
    row_h = 1.35
    for i, rec in enumerate(recommendations, 1):
        # Background alternation
        if i % 2 == 0:
            add_rect(s, 0.6, yy - 0.05, 12.1, row_h, RGBColor(0xFA, 0xFB, 0xFC))

        # Number badge
        num_box = add_rect(s, 0.7, yy + 0.1, 0.35, 0.35, ACCENT)
        add_text(s, str(i), 0.7, yy + 0.1, 0.35, 0.35,
                 size=Pt(11), bold=True, color=BG_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Recommendation
        xx = 0.6
        add_text(s, rec['title'], xx + 0.6, yy + 0.05, 4.7, 0.4,
                 size=Pt(11), bold=True, color=TEXT_PRIMARY)
        add_text(s, rec['detail'], xx + 0.6, yy + 0.45, 4.7, 0.85,
                 size=Pt(10), color=TEXT_SECONDARY)
        xx += 5.4

        # Owner
        add_text(s, rec['owner'], xx + 0.1, yy + 0.15, 2.1, 0.4,
                 size=Pt(11), color=TEXT_PRIMARY)
        add_text(s, rec.get('owner_detail', ''), xx + 0.1, yy + 0.5, 2.1, 0.7,
                 size=Pt(9), color=TEXT_SECONDARY)
        xx += 2.2

        # Timeline
        add_text(s, rec['timeline'], xx + 0.1, yy + 0.15, 1.5, 0.4,
                 size=Pt(11), color=TEXT_PRIMARY)
        xx += 1.6

        # Success metric
        add_text(s, rec['metric'], xx + 0.1, yy + 0.15, 2.8, row_h - 0.3,
                 size=Pt(10), color=TEXT_PRIMARY)

        yy += row_h
        add_line(s, 0.6, yy - 0.05, 12.1, 0, GRID_LIGHT, 0.3)


# ── MAIN ──────────────────────────────────────────────────────────────────────

import os as _os
_DEFAULT_OUT = _os.path.join(
    _os.path.dirname(_os.path.dirname(_os.path.abspath(__file__))),
    "summaries", f"CX_Intelligence_{__import__('datetime').datetime.now().strftime('%Y-%m')}.pptx"
)

def generate_pptx(output_path=None):
    if output_path is None:
        output_path = _DEFAULT_OUT
    configure_matplotlib()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Real data from May 2026 Run #3 ──
    priority_items = [
        ("Delivery & Logistics", 675.9),
        ("App Performance", 261.6),
        ("Checkout & Payment", 243.7),
        ("Return & Refund", 194.4),
        ("Account & Security", 140.4),
        ("Seller Experience", 126.1),
        ("Promotion & Voucher", 72.5),
    ]

    mom_items = [
        ("Return & Refund", 104),
        ("Account & Security", 30),
        ("Delivery & Logistics", 20),
        ("Checkout & Payment", 5),
        ("Seller Experience", -2),
        ("App Performance", -23),
    ]

    # ── SLIDE 1: COVER ──
    build_cover(prs,
        app_name="Shopee Indonesia",
        period="May 6 – 12, 2026",
        run_label="Run #3 · Android + iOS · Automated Pipeline",
        total_reviews=627,
        pct_negative=79.8,
        avg_severity=2.67,
        top_issue="Delivery & Logistics (score 675.9)",
        author="Lidharmawan Suryaatmadja")

    # ── SLIDE 2: EXECUTIVE SUMMARY ──
    build_exec_summary(prs,
        headline=("Delivery is the dominant pain driver — but Return & Refund "
                  "is the emerging churn risk we need to act on this quarter"),
        situation=("627 Shopee Indonesia reviews (520 Android + 107 iOS) from May 6–18 were "
                   "classified and scored by the automated CX Intelligence pipeline, "
                   "with 79.8% negative sentiment and an average severity of 2.67/5."),
        complication=("Three of the six tracked issue areas grew month-over-month, "
                       "and one — Return & Refund — doubled in priority score while "
                       "registering the highest average severity in the dataset (3.70/5)."),
        resolution_items=[
            ("Delivery & Logistics remains the #1 priority (score 675.9, +20% MoM).",
             "Driven by SPX courier protocol breakdown: unauthorized returns after a single missed WA, no phone fallback."),
            ("Return & Refund is the emerging risk (+104% MoM, severity 3.70/5).",
             "Low volume but extreme severity. Qualitative interviews show users abandon the process — actual pain volume is severely undercounted."),
            ("Competitive position: strong on app + checkout, weak on delivery.",
             "Shopee scores 3× better than Tokopedia on Checkout, 3.8× better on Seller Experience — but 2.4× worse on Delivery & Logistics."),
        ],
        page=2, total=7)

    # ── SLIDE 3: HEADLINE KPIS ──
    build_headline_kpis(prs,
        headline=("Sentiment held above 75% negative for the third consecutive "
                  "month, with severity ticking up to 2.67/5"),
        kpis=[
            {"value": "627",       "label": "Reviews analyzed",
             "delta": "520 Android + 107 iOS",  "delta_color": TEXT_SECONDARY},
            {"value": "79.8%",     "label": "Negative sentiment",
             "delta": "▲ +8.8 pp",      "delta_color": SEMANTIC_RED},
            {"value": "2.67 / 5",  "label": "Avg severity",
             "delta": "▲ +0.18",        "delta_color": SEMANTIC_AMBER},
            {"value": "6 / 7",     "label": "Areas with MoM growth",
             "delta": "▲ +2 vs Apr",    "delta_color": SEMANTIC_RED},
        ],
        page=3, total=7)

    # ── SLIDE 4: FINDING #1 — DELIVERY ──
    chart1 = chart_to_image(priority_bar_chart(priority_items,
                                                emphasize_top_n=3,
                                                watch_indices=[3]))
    build_deep_dive(prs,
        headline=("Delivery & Logistics tops the priority list — 150 reviews, "
                  "score 675.9, +20% over Run #1"),
        subhead="The top 3 areas (orange) account for 70% of total priority weight; Return & Refund (amber) is flagged as a watch area.",
        chart_image=chart1,
        quote_text=("masa paket gw datang jam 5 sore tpi sya lupa respon "
                    "kurirnya trus paket gw kena retur"),
        quote_attribution="1★ review · May 8, 2026 · Bekasi",
        root_cause_items=[
            ("SPX courier SOP misalignment",
             "Return triggered after 1 unanswered WA in 2 hours, no phone fallback."),
            ("Hub capacity bottleneck",
             "Bekasi & Lampung corridors show 4-7 day same-city delivery."),
            ("ETA changes post-order without notification",
             "Broken-promise effect drives severity above raw delay impact."),
        ],
        page=4, total=7)

    # ── SLIDE 5: COMPETITIVE — SHOPEE VS TOKOPEDIA ──
    chart2 = chart_to_image(competitive_grouped_bar(
        ["Delivery", "App Perf.", "Checkout", "Seller Exp."],
        [108.5, 65.1, 44.8, 24.7],
        [45.3, 136.0, 134.1, 94.5]))
    build_competitive(prs,
        headline=("Shopee outperforms Tokopedia on app, checkout, and seller "
                  "experience — but trails 2.4× on delivery"),
        subhead="Scores normalized per 100 reviews · Same classifier · Same May 2026 period",
        chart_image=chart2,
        finding_items=[
            ("Shopee's competitive moat",
             "3× better on Checkout, 3.8× better on Seller Experience. Strong UX is the retention floor.",
             ACCENT),
            ("Shopee's exposure",
             "2.4× worse on Delivery. This is the single area where users may rationalize switching to Tokopedia.",
             SEMANTIC_RED),
            ("Hidden trade-off",
             "Tokopedia's better delivery is paid for with 2-3× worse app and checkout. Not an obvious win.",
             SEMANTIC_AMBER),
        ],
        page=5, total=7)

    # ── SLIDE 6: WATCH AREAS ──
    build_watch_area(prs,
        headline=("Two watch areas warrant attention before the June run, "
                  "even though they are not yet in the top 3"),
        subhead="Watch areas are leading indicators — high severity or rapid MoM growth that may not yet show in volume.",
        watch_items=[
            {"name": "Return & Refund — Escalation Risk",
             "signal": ("+104% MoM growth in priority score (95.5 → 194.4). "
                        "Avg severity 3.70/5 is the highest in the dataset."),
             "impact": ("Refund failures cited as reason to switch platforms permanently. "
                        "4/4 interviewed users abandoned the return process mid-flow."),
             "trigger": ("If score > 250 in June run, escalate to dedicated workstream. "
                         "Instrument return funnel completion rate as primary KPI.")},
            {"name": "Account & Security — Steady Climb",
             "signal": ("Priority score up +30% MoM (108.2 → 140.4). "
                        "Severity 3.47/5 is the second highest in dataset."),
             "impact": ("Lost coins, hijack accusations, and credit-score lag erode trust at the foundation. "
                        "Hard to recover once perceived."),
             "trigger": ("If volume crosses 50 reviews/month, treat as critical. "
                         "Monitor coin-loss complaint rate as proxy.")},
        ],
        page=6, total=7)

    # ── SLIDE 7: RECOMMENDATIONS ──
    build_recommendations(prs,
        headline=("Three concrete actions to address May findings, with clear "
                  "owners, timelines, and measurable outcomes"),
        subhead="Recommendations are ordered by expected impact, not by ease of execution.",
        recommendations=[
            {"title": "Enforce 3-attempt courier protocol before any system-initiated return",
             "detail": "Lock the auto-return trigger in the SPX courier app until 3 contact attempts logged (WA + phone + visit).",
             "owner": "Logistics Ops",
             "owner_detail": "+ SPX Partner Team",
             "timeline": "2 weeks",
             "metric": "Unauthorized return complaints down 50% by August run"},
            {"title": "Audit return UX for abandonment friction",
             "detail": "Instrument return-flow funnel: opens → reasons → photo upload → submit. Identify drop-off cliff.",
             "owner": "Returns Product",
             "owner_detail": "+ Data Insights",
             "timeline": "4 weeks",
             "metric": "Return CSAT +0.5 pt; return submission rate measurable by Sept"},
            {"title": "Invest in delivery reliability as Tokopedia closes ground",
             "detail": "Multi-month commitment to address the only category where competitor outperforms us.",
             "owner": "Logistics Strategy",
             "owner_detail": "+ Network Ops",
             "timeline": "Quarterly",
             "metric": "Normalized delivery score per 100 reviews drops below 100"},
        ],
        page=7, total=7)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = generate_pptx()
    print(f"Saved: {path}")
